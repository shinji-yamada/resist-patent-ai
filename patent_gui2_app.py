import streamlit as st
from pptx import Presentation
import os
import io  # ← 重要！

# ページ設定
st.set_page_config(page_title="特許調査ツール", layout="wide")
st.title("🔍 特許調査アプリ（PowerPoint対応）")

# --- ① IPC分類コードのチェックボックス表示 ---
st.header("① IPC分類コードの選択")

ipc_options = {
    "C08F2/00（重合反応）": "C08F2/00",
    "G03F7/027（感光性レジスト）": "G03F7/027",
    "C09D11/00（光学コーティング）": "C09D11/00",
    "C08L33/00（樹脂組成）": "C08L33/00",
}

selected_ipcs = []
st.markdown("調査対象としたいIPC分類を選択してください：")

for label in ipc_options:
    if st.checkbox(label, value=True):
        selected_ipcs.append(ipc_options[label])

# --- ② PowerPointファイルのアップロード ---
st.header("② PowerPoint資料のアップロード（発明概要）")
uploaded_file = st.file_uploader("📎 発明内容を説明するPowerPointファイル（.pptx）をアップロード", type=["pptx"])

ppt_text = ""

if uploaded_file:
    # セッションに保存されていないなら読み取り処理
    if "ppt_text" not in st.session_state:
        try:
            ppt_binary = uploaded_file.read()
            prs = Presentation(io.BytesIO(ppt_binary))
            ppt_text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        ppt_text += shape.text + "\n"
            st.session_state["ppt_text"] = ppt_text  # セッションに保存
            st.success("✅ PowerPointファイルからテキストを抽出しました。")
            st.text_area("📘 抽出されたテキスト", ppt_text, height=200)
        except Exception as e:
            st.error(f"❌ PowerPointの読み込み中にエラーが発生しました: {e}")
    else:
        ppt_text = st.session_state["ppt_text"]
        st.text_area("📘 抽出されたテキスト", ppt_text, height=200)

# --- ③ 実行ボタンと後続処理（Google Patentsジャンプ付き） ---
st.header("③ 類似特許の取得")

if st.button("🔍 類似特許を探す（※現在はダミー動作）"):
    ppt_text = st.session_state.get("ppt_text", "")  # 安全に取り出す

    if not ppt_text:
        st.warning("⚠️ PowerPointファイルから抽出したテキストが見つかりません。")
    elif not selected_ipcs:
        st.warning("⚠️ IPC分類コードを1つ以上選択してください。")
    else:
        st.info("🔧 概念検索機能とGoogle Patents検索は現在開発中です。")
        st.write("✅ 選択されたIPCコード:", ", ".join(selected_ipcs))

        # 🔗 Google Patentsジャンプリンク
        query = "+".join(ppt_text.split()[:10])  # 最初の10単語
        google_patents_url = f"https://patents.google.com/?q={query}"
        st.markdown(f"[🔗 Google Patentsで検索する]({google_patents_url})", unsafe_allow_html=True)

        # 抽出テキスト一部表示
        st.write("📝 抽出されたテキストの一部:", ppt_text[:300] + "..." if len(ppt_text) > 300 else ppt_text)
